"""Create a .pptx deck from a Markdown outline.

Each H2 (##) heading becomes a slide. The heading text is the slide title; the
body (bullets and paragraphs until the next H2) becomes slide content.
An optional title slide (H1 of the document) can be added via --title/--subtitle.
Requires `python-pptx`. Run via the office_python interpreter.
"""

from __future__ import annotations

from pathlib import Path
import argparse
import json
import re


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    markdown = _read_markdown(args)
    if not markdown.strip():
        raise ValueError("Markdown 内容为空。")
    output_path = Path(args.output_path).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    create_pptx(markdown, output_path, title=args.title or "", subtitle=args.subtitle or "")
    slides = _count_slides(markdown)
    print(
        json.dumps(
            {"output_path": str(output_path), "slides": slides},
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create a .pptx deck from a Markdown outline.")
    parser.add_argument("--markdown-path")
    parser.add_argument("--markdown-content")
    parser.add_argument("--output-path", required=True)
    parser.add_argument("--title", default="")
    parser.add_argument("--subtitle", default="")
    return parser.parse_args(argv)


def _read_markdown(args: argparse.Namespace) -> str:
    if args.markdown_path:
        return Path(args.markdown_path).read_text(encoding="utf-8", errors="replace")
    return args.markdown_content or ""


def create_pptx(markdown: str, output_path: Path, *, title: str = "", subtitle: str = "") -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    blank_layout = presentation.slide_layouts[6]

    if title:
        slide = presentation.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.333), Inches(1.0))
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = subtitle
            sub_p.font.size = Pt(22)

    slides = _split_slides(markdown)
    for slide_title, bullets in slides:
        slide = presentation.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = slide_title
        title_p.font.size = Pt(36)
        title_p.font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for index, bullet in enumerate(bullets):
            p = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.level = 0

    presentation.save(str(output_path))


def _split_slides(markdown: str) -> list[tuple[str, list[str]]]:
    """Return [(title, [body lines])] per H2 section."""
    lines = markdown.splitlines()
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_body: list[str] = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("## "):
            if current_title is not None:
                slides.append((current_title, current_body))
            current_title = stripped[3:].strip()
            current_body = []
        elif stripped.startswith("# "):
            # H1 is the document title, skip (handled separately)
            continue
        elif current_title is not None and stripped:
            # Strip leading bullet markers for cleaner body text.
            text = re.sub(r"^\s*([-*•]|\d+\.)\s+", "", stripped)
            current_body.append(text)
    if current_title is not None:
        slides.append((current_title, current_body))
    return slides


def _count_slides(markdown: str) -> int:
    return len(_split_slides(markdown))


if __name__ == "__main__":
    raise SystemExit(main())
