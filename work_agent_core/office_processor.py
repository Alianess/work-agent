from __future__ import annotations

from pathlib import Path
from typing import Any
import argparse
import csv
import json
import os
import subprocess
import sys
import tempfile


TEXT_SEPARATOR = "\n\n"


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    path = Path(args.path).expanduser().resolve()
    if not path.is_file():
        raise FileNotFoundError(path)
    extension = path.suffix.lower()
    if extension == ".docx":
        result = extract_docx(path)
    elif extension in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        result = extract_xlsx(path)
    elif extension in {".csv", ".tsv"}:
        result = extract_delimited(path)
    elif extension == ".pptx":
        result = extract_pptx(path)
    elif extension == ".pdf":
        result = extract_pdf(path)
    else:
        raise ValueError(f"Unsupported office file type: {extension or 'no extension'}")

    markdown = build_markdown(path, result, args.operation)
    if len(markdown) > args.max_chars:
        markdown = markdown[: args.max_chars].rstrip() + "\n\n...[truncated]"
    output_path = Path(args.output_path).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(markdown, encoding="utf-8")
    print(
        json.dumps(
            {
                "source_path": str(path),
                "output_path": str(output_path),
                "kind": result["kind"],
                "metadata": result.get("metadata", {}),
                "excerpt": markdown[:4000],
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


def parse_args(argv: list[str] | None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract office document content to Markdown.")
    parser.add_argument("path")
    parser.add_argument("--operation", default="extract_text")
    parser.add_argument("--output-path", required=True)
    parser.add_argument("--max-chars", type=int, default=50000)
    return parser.parse_args(argv)


def extract_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    document = Document(path)
    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)
    for index, table in enumerate(document.tables, start=1):
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        if rows:
            blocks.append(f"### 表格 {index}\n\n{markdown_table(rows)}")
    props = document.core_properties
    metadata = {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "title": props.title or "",
        "author": props.author or "",
    }
    return {"kind": "docx", "metadata": metadata, "sections": blocks}


def extract_xlsx(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False, read_only=True)
    sections: list[str] = []
    metadata = {"sheets": workbook.sheetnames}
    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        max_rows = min(sheet.max_row or 0, 80)
        max_cols = min(sheet.max_column or 0, 30)
        for row in sheet.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
            values = [format_cell(value) for value in row]
            if any(value for value in values):
                rows.append(values)
        if rows:
            sections.append(f"## 工作表：{sheet.title}\n\n{markdown_table(rows)}")
        else:
            sections.append(f"## 工作表：{sheet.title}\n\n（空表或前 {max_rows} 行无内容）")
    workbook.close()
    return {"kind": "xlsx", "metadata": metadata, "sections": sections}


def extract_delimited(path: Path) -> dict[str, Any]:
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    rows: list[list[str]] = []
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for index, row in enumerate(reader):
            if index >= 100:
                break
            rows.append([cell.strip() for cell in row[:40]])
    return {
        "kind": path.suffix.lower().lstrip("."),
        "metadata": {"rows_sampled": len(rows)},
        "sections": [markdown_table(rows) if rows else "（空文件）"],
    }


def extract_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(path)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        body = "\n\n".join(texts).strip() or "（本页未提取到文本）"
        sections.append(f"## 幻灯片 {index}\n\n{body}")
    metadata = {"slides": len(presentation.slides)}
    return {"kind": "pptx", "metadata": metadata, "sections": sections}


def extract_pdf(path: Path) -> dict[str, Any]:
    poppler_result = extract_pdf_with_poppler(path)
    if poppler_result is not None:
        return poppler_result

    sections: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            metadata["pages"] = len(pdf.pages)
            for index, page in enumerate(pdf.pages[:120], start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    sections.append(f"## 第 {index} 页\n\n{text}")
                tables = page.extract_tables() or []
                for table_index, table in enumerate(tables[:5], start=1):
                    rows = [[format_cell(cell) for cell in row] for row in table if row]
                    if rows:
                        sections.append(f"### 第 {index} 页表格 {table_index}\n\n{markdown_table(rows)}")
    except Exception as error:
        first_error = f"{type(error).__name__}: {error}"
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            metadata["pages"] = len(reader.pages)
            for index, page in enumerate(reader.pages[:120], start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    sections.append(f"## 第 {index} 页\n\n{text}")
        except Exception as fallback_error:
            raise RuntimeError(
                "PDF 提取失败，未生成部分结果，避免误判为已完整读取。"
                f"pdfplumber: {first_error}; pypdf: {type(fallback_error).__name__}: {fallback_error}"
            ) from fallback_error
    if not sections:
        sections.append("（未提取到可读文本；如果这是扫描件，需要后续 OCR 能力。）")
    return {"kind": "pdf", "metadata": metadata, "sections": sections}


def extract_pdf_with_poppler(path: Path) -> dict[str, Any] | None:
    pdftotext = find_runtime_binary("pdftotext")
    if pdftotext is None:
        return None

    metadata = pdf_metadata_with_poppler(path)
    with tempfile.TemporaryDirectory(prefix="work-agent-pdf-") as temp_dir:
        output_path = Path(temp_dir) / "text.txt"
        result = subprocess.run(
            [
                str(pdftotext),
                "-layout",
                "-enc",
                "UTF-8",
                str(path),
                str(output_path),
            ],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=180,
            check=False,
        )
        if result.returncode != 0:
            raise RuntimeError(
                "Poppler pdftotext 读取 PDF 失败，未生成部分结果。"
                f"stderr: {result.stderr.strip() or result.stdout.strip()}"
            )
        text = output_path.read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        return None
    pages = [page.strip() for page in text.split("\f")]
    sections = [
        f"## 第 {index} 页\n\n{page}"
        for index, page in enumerate(pages, start=1)
        if page
    ]
    metadata["extractor"] = "poppler-pdftotext"
    metadata["extracted_pages"] = len(sections)
    return {"kind": "pdf", "metadata": metadata, "sections": sections}


def pdf_metadata_with_poppler(path: Path) -> dict[str, Any]:
    pdfinfo = find_runtime_binary("pdfinfo")
    if pdfinfo is None:
        return {}
    result = subprocess.run(
        [str(pdfinfo), str(path)],
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=60,
        check=False,
    )
    if result.returncode != 0:
        return {"pdfinfo_error": (result.stderr or result.stdout).strip()}
    metadata: dict[str, Any] = {}
    for line in result.stdout.splitlines():
        if ":" not in line:
            continue
        key, value = line.split(":", 1)
        key = key.strip().lower().replace(" ", "_")
        value = value.strip()
        if key in {"title", "subject", "author", "creator", "producer", "pages", "encrypted", "pdf_version"}:
            metadata[key] = value
    return metadata


def find_runtime_binary(name: str) -> Path | None:
    candidates = [
        Path(os.environ.get("WORK_AGENT_RUNTIME_BIN", "")) / name if os.environ.get("WORK_AGENT_RUNTIME_BIN") else None,
        Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/bin/override" / name,
        Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/bin" / name,
        Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/native/poppler/poppler/bin" / name,
        Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/native/poppler/bin" / name,
    ]
    for candidate in candidates:
        if candidate and candidate.is_file():
            return candidate
    return None


def build_markdown(path: Path, result: dict[str, Any], operation: str) -> str:
    metadata = result.get("metadata") or {}
    lines = [
        f"# {path.name}",
        "",
        f"- 文件类型：{result['kind']}",
        f"- 处理方式：{operation}",
    ]
    for key, value in metadata.items():
        lines.append(f"- {key}：{value}")
    sections = [str(section).strip() for section in result.get("sections") or [] if str(section).strip()]
    return "\n".join(lines).strip() + TEXT_SEPARATOR + TEXT_SEPARATOR.join(sections)


def markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = [clean_table_cell(cell) or f"列{i + 1}" for i, cell in enumerate(normalized[0])]
    body = normalized[1:] or [[""] * width]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(clean_table_cell(cell) for cell in row) + " |")
    return "\n".join(lines)


def clean_table_cell(value: Any) -> str:
    return format_cell(value).replace("|", "\\|").replace("\n", " ").strip()


def format_cell(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:
        print(f"office_processor error: {type(error).__name__}: {error}", file=sys.stderr)
        raise SystemExit(1)
